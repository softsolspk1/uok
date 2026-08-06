"""
Builds a PowerPoint presentation summarizing the University of Karachi
website redesign, for presentation to the Vice Chancellor.

Run: python scripts/build_presentation.py
Output: UOK_Website_Presentation.pptx (project root)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------------------------------------------------------------- palette
# Pulled directly from site/index.html CSS custom properties (light theme)
DARK_GREEN   = RGBColor(0x1B, 0x2A, 0x1E)   # --text (deep forest green)
ACCENT_GREEN = RGBColor(0x3B, 0x73, 0x31)   # --accent
GREEN_INV    = RGBColor(0x1E, 0x44, 0x29)   # --bg-inverse
GOLD         = RGBColor(0x84, 0x5E, 0x2F)   # --accent-2
GOLD_SOFT    = RGBColor(0xC9, 0xA5, 0x67)   # --accent-2-soft
BEIGE        = RGBColor(0xED, 0xE6, 0xD2)   # --bg
CREAM        = RGBColor(0xF7, 0xF3, 0xE4)   # --bg-elevated
WHITE_TXT    = RGBColor(0xED, 0xE8, 0xD6)   # --text on dark
MUTED        = RGBColor(0x4B, 0x5A, 0x45)   # --text-soft

FONT_HEAD = "Cambria"      # stand-in for Fraunces (display serif)
FONT_BODY = "Calibri"      # stand-in for Work Sans

SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        no_line(shp)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def set_run(run, text, size, color, font=FONT_BODY, bold=False, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic


def footer(slide, page_no, label="UNIVERSITY OF KARACHI — WEBSITE REDESIGN"):
    rect(slide, Inches(0.55), Inches(7.08), Inches(12.23), Pt(1.1), GOLD_SOFT)
    _, tf = textbox(slide, Inches(0.55), Inches(7.12), Inches(9), Inches(0.3))
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, label, 9, MUTED, font=FONT_BODY, bold=True)
    _, tf2 = textbox(slide, Inches(11.5), Inches(7.12), Inches(1.28), Inches(0.3))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    set_run(r2, str(page_no), 9, MUTED, bold=True)


def header(slide, kicker, title):
    rect(slide, 0, 0, SW, Inches(1.5), GREEN_INV)
    rect(slide, 0, Inches(1.5), SW, Pt(3.5), GOLD)
    _, tf = textbox(slide, Inches(0.55), Inches(0.18), Inches(11), Inches(0.35))
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, kicker.upper(), 12, GOLD_SOFT, bold=True)
    p.runs[0].font._rPr.set('spc', '200')
    _, tf2 = textbox(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.95))
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    set_run(r2, title, 30, WHITE_TXT, font=FONT_HEAD, bold=True)


def bullets(slide, x, y, w, h, items, size=15, color=DARK_GREEN,
            marker_color=GOLD, gap=10, bold_lead=True):
    _, tf = textbox(slide, x, y, w, h)
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = 1.12
        rm = p.add_run()
        set_run(rm, "■  ", size, marker_color, bold=True)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            set_run(r1, lead, size, color, bold=True)
            r2 = p.add_run()
            set_run(r2, rest, size, color)
        else:
            r1 = p.add_run()
            set_run(r1, item, size, color)
    return tf


def bullet_block_y(n_items, size, gap, top=1.75, bottom=6.9):
    """Vertical start position that centers an n-item bullet block
    within the slide's content area, instead of always hugging the
    header and leaving dead space at the bottom."""
    line_in = size * 1.12 / 72.0
    block_in = n_items * line_in + max(n_items - 1, 0) * (gap / 72.0)
    avail = bottom - top
    y = top + max(0.0, (avail - block_in) / 2.0)
    return Inches(y)


def content_slide(page_no, kicker, title, bg=BEIGE):
    s = add_slide()
    set_bg(s, bg)
    header(s, kicker, title)
    footer(s, page_no)
    return s


# ============================================================ 1. TITLE
s = add_slide()
set_bg(s, GREEN_INV)
rect(s, 0, Inches(6.55), SW, Inches(0.95), DARK_GREEN)
rect(s, Inches(0.9), Inches(2.55), Inches(1.4), Pt(3.5), GOLD)

_, tf = textbox(s, Inches(0.9), Inches(0.75), Inches(11.5), Inches(0.5))
p = tf.paragraphs[0]
r = p.add_run()
set_run(r, "UNIVERSITY OF KARACHI  ·  EST. 1951", 14, GOLD_SOFT, bold=True)

_, tf = textbox(s, Inches(0.85), Inches(2.75), Inches(11.6), Inches(2.2))
p = tf.paragraphs[0]
r = p.add_run()
set_run(r, "The New University Website", 44, WHITE_TXT, font=FONT_HEAD, bold=True)
p2 = tf.add_paragraph()
p2.space_before = Pt(6)
r2 = p2.add_run()
set_run(r2, "A Complete Digital Redesign", 44, WHITE_TXT, font=FONT_HEAD, bold=True)

_, tf = textbox(s, Inches(0.9), Inches(5.05), Inches(10.5), Inches(0.6))
p = tf.paragraphs[0]
r = p.add_run()
set_run(r, "Prepared for presentation to the Vice Chancellor", 18, BEIGE, italic=True)

_, tf = textbox(s, Inches(0.9), Inches(6.75), Inches(9), Inches(0.5))
p = tf.paragraphs[0]
r = p.add_run()
set_run(r, "Web Development Team  |  August 2026", 13, GOLD_SOFT, bold=True)


# ============================================================ 2. AGENDA
s = content_slide(2, "Overview", "What We'll Cover")
items = [
    "Project overview & goals",
    "Design philosophy — colour, typography, identity",
    "Homepage walkthrough",
    "Site-wide features & navigation",
    "Accessibility & responsiveness",
    "Content scope — every page built",
    "Template-driven architecture",
    "Complete feature list",
    "Technical foundation",
    "Benefits to the University",
    "Next steps to go live",
]
col_w = Inches(5.7)
half = (len(items) + 1) // 2
agenda_y = bullet_block_y(half, 17, 16)
bullets(s, Inches(0.7), agenda_y, col_w, Inches(5), items[:half], size=17, gap=16)
bullets(s, Inches(6.7), agenda_y, col_w, Inches(5), items[half:], size=17, gap=16)


# ============================================================ 3. PROJECT OVERVIEW
s = content_slide(3, "Project Overview", "Reimagining KU's Front Door Online")
_, tf = textbox(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(1.1))
p = tf.paragraphs[0]
r = p.add_run()
set_run(r, "The official website has been rebuilt from the ground up — every page, "
           "every section — to give the University a digital presence that matches "
           "its 75-year academic legacy.", 16, DARK_GREEN)

stat_data = [("96", "Pages Built"), ("57", "Department Pages"),
             ("13", "Institutes & Centres"), ("8", "Faculties Represented")]
x = Inches(0.7)
for label_val, label_name in stat_data:
    card = rect(s, x, Inches(3.15), Inches(2.75), Inches(1.7), CREAM)
    tf = card.text_frame
    tf.margin_left = tf.margin_right = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, label_val, 34, ACCENT_GREEN, font=FONT_HEAD, bold=True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    set_run(r2, label_name.upper(), 11, GOLD, bold=True)
    x += Inches(3.0)

bullets(s, Inches(0.7), Inches(5.25), Inches(11.9), Inches(1.6), [
    ("Scope:  ", "Homepage plus every faculty, department, institute, admissions, "
                 "library, research, examinations and student-services page."),
    ("Approach:  ", "A single consistent design system applied site-wide, built and "
                     "maintained through reusable page templates."),
], size=14, gap=8)


# ============================================================ 4. DESIGN PHILOSOPHY
s = content_slide(4, "Design Philosophy", "A Visual Identity Rooted in Heritage")
_, tf = textbox(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(0.6))
p = tf.paragraphs[0]
r = p.add_run()
set_run(r, "Deep academic green, warm parchment beige, and gold accents — "
           "a palette that feels institutional, credible and timeless.", 15, DARK_GREEN, italic=True)

swatches = [("#EDE6D2", "Parchment Beige", BEIGE), ("#1B2A1E", "Academic Green", DARK_GREEN),
            ("#3B7331", "Accent Green", ACCENT_GREEN), ("#845E2F", "Heritage Gold", GOLD)]
x = Inches(0.7)
for hexcode, name, color in swatches:
    swatch = rect(s, x, Inches(2.6), Inches(2.75), Inches(1.0), color, line=True)
    swatch.line.color.rgb = MUTED
    swatch.line.width = Pt(1)
    _, tf = textbox(s, x, Inches(3.65), Inches(2.75), Inches(0.7))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, name, 12, DARK_GREEN, bold=True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    set_run(r2, hexcode, 10, MUTED)
    x += Inches(3.0)

bullets(s, Inches(0.7), Inches(4.75), Inches(11.9), Inches(2), [
    ("Typography:  ", "Fraunces (serif display) for gravitas, Work Sans for clean "
                       "readability, and Noto Nastaliq Urdu for the bilingual identity."),
    ("Dark mode:  ", "An automatic dark theme is built in, respecting each visitor's "
                      "system preference."),
], size=14, gap=8)


# ============================================================ 5. HOMEPAGE WALKTHROUGH
s = content_slide(5, "Homepage", "A Homepage That Tells KU's Story")
items = [
    ("Hero section — ", "bilingual English/Urdu welcome with quick-task shortcuts"),
    ("Stats band — ", "key institutional numbers at a glance"),
    ("History timeline — ", "“Established by an act of parliament, seventy-five years ago”"),
    ("Academics — ", "“Eight faculties. Fifty-three departments.” directory"),
    ("Research spotlight — ", "“Where Pakistani chemistry earned its global reputation”"),
    ("Library feature — ", "Dr. Mahmud Hussain Library, with key facts highlighted"),
    ("News & notices — ", "tabbed feed for Notifications, Examinations and Campus Life"),
]
bullets(s, Inches(0.7), bullet_block_y(len(items), 15.5, 13), Inches(11.9), Inches(5), items, size=15.5, gap=13)


# ============================================================ 6. SITE-WIDE NAVIGATION
s = content_slide(6, "Navigation", "Consistent, Intuitive Navigation Everywhere")
items = [
    ("Dropdown menus — ", "structured mega-menus across every one of the 96 pages"),
    ("Mobile menu — ", "accessible hamburger navigation for phones and tablets"),
    ("Bilingual toggle — ", "instant English / Urdu language switch"),
    ("Consistent branding — ", "logo, Home button and header identical on every page"),
    ("Unified footer — ", "quick links, contact details and credits on every page"),
]
bullets(s, Inches(0.7), bullet_block_y(len(items), 16, 18), Inches(11.9), Inches(5), items, size=16, gap=18)


# ============================================================ 7. ACCESSIBILITY
s = content_slide(7, "Accessibility", "Built for Every Visitor")
items = [
    ("WCAG AA contrast — ", "verified text/background contrast across all pages"),
    ("Screen-reader support — ", "ARIA labels, roles and states throughout "
                                  "(aria-expanded, aria-pressed, role=\"tablist\")"),
    ("Reduced-motion support — ", "honours visitors' system motion preferences"),
    ("Fully responsive — ", "tuned breakpoints for mobile, tablet and desktop"),
    ("Light & dark themes — ", "automatic, based on system preference"),
]
bullets(s, Inches(0.7), bullet_block_y(len(items), 16, 18), Inches(11.9), Inches(5), items, size=16, gap=18)


# ============================================================ 8. CONTENT SCOPE
s = content_slide(8, "Content Scope", "Every Corner of the University, Online")
rows = [
    ("57", "Department pages — one for every academic department"),
    ("13", "Institute & research-centre pages"),
    ("8", "Faculties overview"),
    ("1", "Admissions & Prospectus"),
    ("1", "Examinations portal"),
    ("1", "Library & Research pages"),
    ("1", "Student Portal, Alumni & Journals"),
    ("1", "Administration & Contact"),
]
y = Inches(1.85)
col = 0
xs = [Inches(0.7), Inches(6.9)]
for i, (num, desc) in enumerate(rows):
    col = i % 2
    row_y = y + Inches((i // 2) * 1.15)
    card = rect(s, xs[col], row_y, Inches(5.7), Inches(1.0), CREAM)
    tf = card.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, num + "   ", 22, ACCENT_GREEN, font=FONT_HEAD, bold=True)
    r2 = p.add_run()
    set_run(r2, desc, 13.5, DARK_GREEN)


# ============================================================ 9. TEMPLATE ARCHITECTURE
s = content_slide(9, "Architecture", "One System, Ninety-Six Pages")
_, tf = textbox(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(1.0))
p = tf.paragraphs[0]
r = p.add_run()
set_run(r, "Every department and institute page is generated from a shared template, "
           "so all pages stay visually and structurally consistent.", 16, DARK_GREEN)
items = [
    ("Shared templates — ", "department_template, institute, admissions, library, "
                             "research and student-portal templates"),
    ("Automated build scripts — ", "Python scripts assemble each page from the "
                                    "template, ensuring zero drift between pages"),
    ("One update, everywhere — ", "a change to the shared template can be rolled "
                                   "out across all 57 department pages and 13 institutes"),
    ("Faster future growth — ", "new departments or centres can be added in minutes"),
]
bullets(s, Inches(0.7), bullet_block_y(len(items), 15.5, 14, top=2.9, bottom=6.9), Inches(11.9), Inches(4), items, size=15.5, gap=14)


# ============================================================ 10. FULL FEATURE LIST (2 slides)
def feature_list_slide(page_no, title, features):
    s = content_slide(page_no, "Feature Summary", title)
    half = (len(features) + 1) // 2
    col_y = bullet_block_y(half, 13.5, 11)
    bullets(s, Inches(0.7), col_y, Inches(5.9), Inches(5), features[:half], size=13.5, gap=11)
    bullets(s, Inches(6.75), col_y, Inches(5.9), Inches(5), features[half:], size=13.5, gap=11)
    return s

features_a = [
    "Bilingual English / Urdu interface",
    "Dropdown navigation menus site-wide",
    "Mobile-responsive hamburger menu",
    "Automatic light & dark mode",
    "WCAG AA accessibility compliance",
    "Custom heritage typography system",
    "Bilingual hero section with quick-task shortcuts",
    "Live institutional stats band",
    "Interactive university history timeline",
    "Academics directory (8 faculties, 53 departments)",
    "Research spotlight section",
    "Library feature with key facts",
]
feature_list_slide(10, "Complete Feature List (1 of 2)", features_a)

features_b = [
    "Tabbed news & notices feed",
    "57 individual department pages",
    "13 institute / research-centre pages",
    "Admissions & prospectus pages",
    "Examinations portal",
    "Student portal page",
    "Alumni page",
    "Administration directory",
    "Contact page",
    "Consistent header, footer & branding on every page",
    "Custom favicon & site identity",
    "Clean URL routing via Vercel",
    "Reduced-motion support for accessibility",
]
feature_list_slide(11, "Complete Feature List (2 of 2)", features_b)


# ============================================================ 12. TECHNICAL FOUNDATION
s = content_slide(12, "Technical Foundation", "Fast, Lightweight, Maintainable")
items = [
    ("Lightweight build — ", "static HTML, CSS and JavaScript with no heavy framework "
                              "overhead, meaning fast page loads"),
    ("Vercel deployment — ", "clean, memorable URLs (e.g. /admissions) via routing rules"),
    ("Template-driven pages — ", "Python assembly scripts keep all 96 pages consistent "
                                  "and easy to regenerate"),
    ("Scalable architecture — ", "built to grow as the University adds departments, "
                                  "programmes and content"),
]
bullets(s, Inches(0.7), bullet_block_y(len(items), 16, 18), Inches(11.9), Inches(5), items, size=16, gap=18)


# ============================================================ 13. BENEFITS
s = content_slide(13, "Benefits", "What This Means for the University")
items = [
    ("Stronger first impression — ", "a modern, credible site for prospective students, "
                                      "faculty and international partners"),
    ("Wider reach — ", "bilingual and accessible design serves more of Pakistan's students"),
    ("Mobile-first — ", "built for the majority of visitors who browse on their phones"),
    ("Lower maintenance cost — ", "templated architecture means fast, low-risk updates"),
    ("Future-ready — ", "a foundation that scales as the University's digital needs grow"),
]
bullets(s, Inches(0.7), bullet_block_y(len(items), 16, 16), Inches(11.9), Inches(5), items, size=16, gap=16)


# ============================================================ 14. NEXT STEPS
s = content_slide(14, "Next Steps", "Path to Going Live")
items = [
    "Content review and sign-off by faculties, departments and administration",
    "Populate live data feeds — notices, results, admissions dates",
    "Cross-browser and cross-device quality assurance",
    "Staff training for ongoing content updates",
    "Domain migration and official go-live",
]
bullets(s, Inches(0.7), bullet_block_y(len(items), 17, 20), Inches(11.9), Inches(5), items, size=17, gap=20)


# ============================================================ 15. THANK YOU
s = add_slide()
set_bg(s, GREEN_INV)
rect(s, Inches(0.9), Inches(3.15), Inches(1.4), Pt(3.5), GOLD)
_, tf = textbox(s, Inches(0.9), Inches(2.4), Inches(11), Inches(1.0))
p = tf.paragraphs[0]
r = p.add_run()
set_run(r, "Thank You", 44, WHITE_TXT, font=FONT_HEAD, bold=True)
_, tf = textbox(s, Inches(0.9), Inches(3.45), Inches(10.5), Inches(0.6))
p = tf.paragraphs[0]
r = p.add_run()
set_run(r, "Questions & Discussion", 18, BEIGE, italic=True)
_, tf = textbox(s, Inches(0.9), Inches(6.6), Inches(10.5), Inches(0.5))
p = tf.paragraphs[0]
r = p.add_run()
set_run(r, "University of Karachi  ·  Web Development Team", 13, GOLD_SOFT, bold=True)


# ---------------------------------------------------------------- save
out_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                         "UOK_Website_Presentation.pptx")
prs.save(out_path)
print("Saved:", out_path)
