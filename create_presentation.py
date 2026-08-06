import sys
import subprocess

# Auto-install python-pptx if it's not installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # ---------------------------------------------------------
    # Helper: Set slide background color (Dark Mode Theme)
    # ---------------------------------------------------------
    def set_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(27, 42, 30) # UoK Dark Green Theme

    # ---------------------------------------------------------
    # Slide 0: Title Slide
    # ---------------------------------------------------------
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    set_dark_background(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Revitalizing the Digital Identity of University of Karachi"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "A Modern Web Experience\nPresented to the Vice Chancellor"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(201, 165, 103) # Gold Accent
    subtitle.text_frame.paragraphs[0].font.italic = True
    if len(subtitle.text_frame.paragraphs) > 1:
        subtitle.text_frame.paragraphs[1].font.color.rgb = RGBColor(201, 165, 103)

    # ---------------------------------------------------------
    # Slide 1: The Vision
    # ---------------------------------------------------------
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    set_dark_background(slide)
    
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "1. The Vision"
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    tf = body_shape.text_frame
    tf.text = "A web presence that reflects the prestige of UoK (Est. 1951)"
    tf.paragraphs[0].font.color.rgb = RGBColor(237, 232, 214)
    
    p = tf.add_paragraph()
    p.text = "Transitioning from legacy infrastructure to a modern, dynamic framework."
    p.font.color.rgb = RGBColor(237, 232, 214)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Prioritizing accessibility, aesthetic excellence, and institutional pride."
    p.font.color.rgb = RGBColor(237, 232, 214)
    p.level = 1

    # ---------------------------------------------------------
    # Slide 2: Modern Design Aesthetics
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    set_dark_background(slide)
    
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    
    title_shape.text = "2. Modern Design Aesthetics"
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    tf = body_shape.text_frame
    tf.text = "Visual Excellence & User Engagement"
    tf.paragraphs[0].font.color.rgb = RGBColor(237, 232, 214)
    
    p = tf.add_paragraph()
    p.text = "Sleek typography (Fraunces & Work Sans) improving readability."
    p.font.color.rgb = RGBColor(237, 232, 214)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Automatic Dark Mode support adapting to user system preferences seamlessly."
    p.font.color.rgb = RGBColor(237, 232, 214)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Micro-animations and engaging hover effects creating a 'live' feel."
    p.font.color.rgb = RGBColor(237, 232, 214)
    p.level = 1

    # ---------------------------------------------------------
    # Slide 3: Enhanced Navigation & Mobile
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    set_dark_background(slide)
    
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    
    title_shape.text = "3. Navigation & Mobile-First Approach"
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    tf = body_shape.text_frame
    tf.text = "Intuitive architecture for over 100+ pages"
    tf.paragraphs[0].font.color.rgb = RGBColor(237, 232, 214)
    
    p = tf.add_paragraph()
    p.text = "Brand new Global Dropdown Menus ensuring 70+ departments are easily reachable."
    p.font.color.rgb = RGBColor(237, 232, 214)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Mobile-First Design: The entire UI reflows flawlessly for smartphones."
    p.font.color.rgb = RGBColor(237, 232, 214)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Custom touch-optimized mobile navigation preventing layout breaking."
    p.font.color.rgb = RGBColor(237, 232, 214)
    p.level = 1

    # ---------------------------------------------------------
    # Slide 4: Automated & Scalable Architecture
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    set_dark_background(slide)
    
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    
    title_shape.text = "4. Automated & Scalable Architecture"
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    tf = body_shape.text_frame
    tf.text = "Built for effortless maintenance and rapid content deployment"
    tf.paragraphs[0].font.color.rgb = RGBColor(237, 232, 214)
    
    p = tf.add_paragraph()
    p.text = "Python-driven generation system (gen_departments.py) for all subpages."
    p.font.color.rgb = RGBColor(237, 232, 214)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Centralized templates ensuring brand consistency across every single page."
    p.font.color.rgb = RGBColor(237, 232, 214)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Instant global updates: Changing a menu or layout instantly applies to 100+ pages."
    p.font.color.rgb = RGBColor(237, 232, 214)
    p.level = 1

    # ---------------------------------------------------------
    # Slide 5: Thank You & Live Demo
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(title_slide_layout)
    set_dark_background(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    subtitle.text = "Questions & Live Demo"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(201, 165, 103)

    # Save the presentation
    output_file = 'UoK_Website_Demo_Presentation.pptx'
    prs.save(output_file)
    print(f"Successfully generated {output_file}!")

if __name__ == '__main__':
    create_presentation()
